import os
import traceback
from pathlib import Path
from typing import Dict, Any, List

# --- Safe External Library Imports ---
try:
    import fitz
except ImportError:
    fitz = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from pdf2docx import Converter as PDFConverter
except ImportError:
    PDFConverter = None

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
except ImportError:
    canvas = None
    letter = None


async def convert_file(
    input_path: str,
    output_path: str,
    source_format: str,
    target_format: str,
) -> Dict[str, Any]:
    """
    Main conversion engine that delegates to specific converters
    based on source and target formats.
    Returns a dictionary with 'success': True/False and 'error' message if failed.
    """
    try:
        src: str = source_format.lower()
        tgt: str = target_format.lower()

        # --- PDF to WORD (DOCX) ---
        if src == "pdf" and tgt == "docx":
            if not PDFConverter:
                return {"success": False, "error": "pdf2docx library not installed"}
            cv = PDFConverter(input_path)
            try:
                cv.convert(output_path, start=0, end=None)
            finally:
                cv.close()
            return {"success": True}

        # --- WORD (DOCX) to PDF ---
        elif src == "docx" and tgt == "pdf":
            if not Document or not canvas:
                return {"success": False, "error": "docx or reportlab library not installed"}
            doc_file = Document(input_path)
            c = canvas.Canvas(output_path, pagesize=letter)
            height: float = letter[1]
            y: float = height - 50
            for para in doc_file.paragraphs:
                if y < 50:
                    c.showPage()
                    y = height - 50
                c.drawString(50, y, para.text[:100])
                y -= 15
            c.save()
            return {"success": True}

        # --- PDF to PPTX ---
        elif src == "pdf" and tgt == "pptx":
            if not Presentation or not fitz:
                return {"success": False, "error": "pptx or fitz library not installed"}
            prs = Presentation()
            with fitz.open(input_path) as doc_pdf:
                for page_num in range(len(doc_pdf)):
                    page = doc_pdf.load_page(page_num)
                    pix = page.get_pixmap()
                    temp_img: str = f"{output_path}_temp_{page_num}.png"
                    pix.save(temp_img)
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(
                        temp_img, 0, 0,
                        width=prs.slide_width,
                        height=prs.slide_height,
                    )
                    os.remove(temp_img)
            prs.save(output_path)
            return {"success": True}

        # --- PDF to IMAGE (PNG/JPG) ---
        elif src == "pdf" and tgt in ("png", "jpg", "jpeg"):
            if not fitz:
                return {"success": False, "error": "fitz library not installed"}
            with fitz.open(input_path) as doc_pdf:
                if len(doc_pdf) == 0:
                    return {"success": False, "error": "PDF is empty"}
                page = doc_pdf.load_page(0)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                pix.save(output_path)
            return {"success": True}

        # --- PDF to EXCEL (XLSX) ---
        elif src == "pdf" and tgt == "xlsx":
            if not fitz or not pd:
                return {"success": False, "error": "fitz or pandas library not installed"}
            with fitz.open(input_path) as doc_pdf:
                text_data: List[List[str]] = []
                for page in doc_pdf:
                    lines: List[str] = page.get_text().split("\n")
                    text_data.extend(
                        [line.split() for line in lines if line.strip()]
                    )
            df = pd.DataFrame(text_data)
            df.to_excel(output_path, index=False, header=False)
            return {"success": True}

        # --- EXCEL (XLSX) to PDF ---
        elif src == "xlsx" and tgt == "pdf":
            if not pd or not canvas:
                return {"success": False, "error": "pandas or reportlab library not installed"}
            df = pd.read_excel(input_path)
            c = canvas.Canvas(output_path, pagesize=letter)
            height = letter[1]
            y = height - 50
            c.drawString(50, y, f"Excel Data Export: {Path(input_path).name}")
            y -= 30
            for _, row in df.iterrows():
                if y < 50:
                    c.showPage()
                    y = height - 50
                row_str: str = " | ".join([str(val) for val in row.values])
                c.drawString(50, y, row_str[:120])
                y -= 15
            c.save()
            return {"success": True}

        # --- PPTX to PDF ---
        elif src == "pptx" and tgt == "pdf":
            if not Presentation or not canvas:
                return {"success": False, "error": "pptx or reportlab library not installed"}
            prs = Presentation(input_path)
            c = canvas.Canvas(output_path, pagesize=letter)
            height = letter[1]
            for i, slide in enumerate(prs.slides):
                y = height - 50
                c.drawString(50, y, f"Slide {i + 1}")
                y -= 30
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if y < 50:
                            c.showPage()
                            y = height - 50
                        c.drawString(50, y, shape.text[:100])
                        y -= 15
                c.showPage()
            c.save()
            return {"success": True}

        # --- PDF to HTML / SVG / TXT ---
        elif src == "pdf" and tgt in ("html", "svg", "txt"):
            if not fitz:
                return {"success": False, "error": "fitz library not installed"}
            content: str = ""
            with fitz.open(input_path) as doc_pdf:
                if tgt == "html":
                    content = "".join(
                        [page.get_text("html") for page in doc_pdf]
                    )
                elif tgt == "svg":
                    if len(doc_pdf) == 0:
                        return {"success": False, "error": "PDF is empty"}
                    content = doc_pdf[0].get_svg_image()
                else:
                    content = "".join(
                        [page.get_text() for page in doc_pdf]
                    )
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(content)
            return {"success": True}

        # --- IMAGE to IMAGE ---
        elif (
            src in ("jpg", "jpeg", "png", "webp", "bmp", "tiff")
            and tgt in ("jpg", "jpeg", "png", "webp", "bmp", "tiff")
        ):
            if not Image:
                return {"success": False, "error": "PIL library not installed"}
            with Image.open(input_path) as img:
                if tgt in ("jpg", "jpeg") and img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                img.save(output_path)
            return {"success": True}

        # --- IMAGE (JPG/PNG) to PDF ---
        elif src in ("jpg", "jpeg", "png") and tgt == "pdf":
            if not fitz:
                return {"success": False, "error": "fitz library not installed"}
            doc_new = fitz.open()
            img_doc = fitz.open(input_path)
            pdf_bytes = img_doc.convert_to_pdf()
            img_doc.close()
            img_pdf = fitz.open("pdf", pdf_bytes)
            doc_new.insert_pdf(img_pdf)
            img_pdf.close()
            doc_new.save(output_path)
            doc_new.close()
            return {"success": True}

        # --- TXT to PDF ---
        elif src == "txt" and tgt == "pdf":
            if not canvas:
                return {"success": False, "error": "reportlab library not installed"}
            with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
                lines = f.readlines()
            c = canvas.Canvas(output_path, pagesize=letter)
            height = letter[1]
            y = height - 50
            for line in lines:
                if y < 50:
                    c.showPage()
                    y = height - 50
                c.drawString(50, y, line.strip()[:100])
                y -= 15
            c.save()
            return {"success": True}

        else:
            return {"success": False, "error": f"Unsupported: {src} to {tgt}"}

    except Exception as e:
        error_msg = f"Conversion error ({source_format} to {target_format}): {str(e)}"
        print(error_msg)
        traceback.print_exc()
        return {"success": False, "error": error_msg}
